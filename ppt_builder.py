"""
Main entry point: assemble the Sunday worship PPTX from agenda + inputs.
"""

import os
import shutil
from pptx import Presentation

from slide_finder import find_slide
from slide_copier import copy_slide, clear_slides, get_layout
from slide_planner import plan_slides
from slide_generators import (
    add_placeholder_slide, add_anthem_title_slide, add_scripture_title_slide,
    add_scripture_verse_slide, add_lyrics_slide, add_sermon_title_slide,
    add_sermon_point_slide, add_announcement_slide
)


def build_pptx(template_path, agenda, input_files, output_path, library_paths=None, intro_path=None, bible_page=None):
    """
    Build the worship PPTX.

    Args:
        template_path: Path to template PPTX (for slide master/layouts)
        agenda: Parsed agenda dict from pdf_parser
        input_files: Dict of filename → filepath for input files
        output_path: Where to save the result
        library_paths: Additional PPTX files to search for hymn/reading slides
        intro_path: Optional fixed-intro PPTX whose slides replace first few
        bible_page: Page number for scripture title (from user)

    Returns:
        Path to the created PPTX
    """
    template = Presentation(template_path)
    use_intro = intro_path and os.path.exists(intro_path)

    # Build library: template + any extra library files
    libraries = [template]
    for p in (library_paths or []):
        if p != template_path and os.path.exists(p):
            try:
                libraries.append(Presentation(p))
            except Exception:
                pass

    # Plan all slides
    slides_to_add = plan_slides(
        template, libraries, agenda, input_files,
        skip_intro=use_intro, bible_page=bible_page
    )

    # Choose output base: intro PPTX or template
    base_path = intro_path if use_intro else template_path
    shutil.copy2(base_path, output_path)
    out_prs = Presentation(output_path)

    if not use_intro:
        clear_slides(out_prs)

    # Add all slides
    for spec in slides_to_add:
        _add_slide(out_prs, template, spec, agenda)

    out_prs.save(output_path)
    return output_path


def _add_slide(out_prs, template, spec, agenda):
    """Dispatch slide creation based on spec type."""
    stype = spec["type"]

    if stype == "copy_template":
        prs = spec.get("prs", template)
        index = spec.get("index")
        if index is not None:
            copy_slide(out_prs, prs, index)

    elif stype == "copy_external":
        prs = spec.get("prs")
        index = spec.get("index")
        if index is not None and prs:
            copy_slide(out_prs, prs, index)

    elif stype == "blank":
        layout = get_layout(out_prs, "Blank")
        out_prs.slides.add_slide(layout)

    elif stype == "hymn_placeholder":
        add_placeholder_slide(out_prs, spec.get("label", ""))

    elif stype == "anthem_title":
        add_anthem_title_slide(out_prs, spec.get("title", ""))

    elif stype == "anthem_lyrics":
        add_lyrics_slide(out_prs, spec.get("title", ""), spec.get("lyrics", ""))

    elif stype == "scripture_title":
        add_scripture_title_slide(out_prs, spec.get("item", {}), spec.get("bible_page"))

    elif stype == "scripture_verses":
        add_scripture_verse_slide(out_prs, spec.get("ref", ""), spec.get("verses", []))

    elif stype == "sermon_title":
        add_sermon_title_slide(out_prs, spec)

    elif stype == "sermon_point":
        add_sermon_point_slide(out_prs, spec)

    elif stype == "announcement":
        add_announcement_slide(out_prs, spec)
