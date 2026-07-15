"""
Low-level XML manipulation for copying slides from template/library PPTXs.
"""

import copy
from lxml import etree
from pptx.oxml.ns import qn


def copy_slide(out_prs, src_prs, index):
    """
    Copy a slide from src_prs at index into out_prs.

    Performs deep XML copy of all shapes and background, preserving formatting
    and images. Relationships (embedded media, hyperlinks) are NOT copied due
    to python-pptx limitations.

    Args:
        out_prs: output Presentation
        src_prs: source Presentation
        index: 0-based slide index in src_prs

    Returns:
        The newly added slide, or None if index is out of range.
    """
    if index is None or index >= len(src_prs.slides):
        return None

    src_slide = src_prs.slides[index]

    # Use a matching layout if possible
    layout_name = src_slide.slide_layout.name
    layout = get_layout(out_prs, layout_name)

    new_slide = out_prs.slides.add_slide(layout)

    # Replace shape tree content
    src_sp_tree = src_slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    new_sp_tree = new_slide._element.find(qn("p:cSld")).find(qn("p:spTree"))

    # Remove auto-generated placeholders
    for child in list(new_sp_tree):
        new_sp_tree.remove(child)

    # Copy shapes from source
    for child in src_sp_tree:
        new_sp_tree.append(copy.deepcopy(child))

    # Copy background
    src_cSld = src_slide._element.find(qn("p:cSld"))
    new_cSld = new_slide._element.find(qn("p:cSld"))
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        new_bg = new_cSld.find(qn("p:bg"))
        if new_bg is not None:
            new_cSld.remove(new_bg)
        new_cSld.insert(0, copy.deepcopy(src_bg))

    return new_slide


def restyle_responsive_slide(slide, slide_width):
    """
    Restyle a copied 啟應文 (responsive reading) slide to match the reference:

    - Title box (contains 啟應文): horizontally centered at the top.
    - Body box: paragraphs LEFT-aligned, manual line-wrap continuations merged
      back into their parent paragraph, box widened to the slide edge.
    """
    from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if not text.strip():
            continue

        if "啟應文" in text:
            # Title: center horizontally, pin to top (skip full-width titles)
            if shape.width < slide_width:
                shape.left = (slide_width - shape.width) // 2
                shape.top = 0
        else:
            # Body: merge continuation lines, left-align, widen to slide edge
            _merge_continuation_paragraphs(shape.text_frame)
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            shape.width = slide_width - max(0, shape.left)


def _merge_continuation_paragraphs(text_frame):
    """
    Merge paragraphs that start with whitespace (manual line-wrap continuations
    from narrow source boxes) into the preceding paragraph.
    """
    paragraphs = list(text_frame.paragraphs)
    for p in paragraphs[1:]:
        if not p.runs or not p.text or not p.text[0].isspace():
            continue
        prev = p._p.getprevious()
        if prev is None or not prev.tag == qn("a:p"):
            continue

        # Drop leading whitespace-only runs, lstrip the first content run
        runs = list(p.runs)
        while runs and not runs[0].text.strip():
            runs[0]._r.getparent().remove(runs[0]._r)
            runs.pop(0)
        if runs:
            runs[0].text = runs[0].text.lstrip()

        # Move remaining runs into the previous paragraph
        for r in runs:
            prev.append(r._r)
        p._p.getparent().remove(p._p)


def clear_slides(prs):
    """Remove all slides from a presentation (no undo)."""
    sldIdLst = prs.slides._sldIdLst
    for i in range(len(prs.slides) - 1, -1, -1):
        sld_id = sldIdLst[i]
        rId = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sld_id)


def get_layout(prs, name):
    """
    Get a slide layout by name from prs.

    Returns the matching layout, or the first layout if not found.
    """
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[0]
