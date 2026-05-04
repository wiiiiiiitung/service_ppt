"""
Slide planning logic: translate agenda + input files → ordered list of slide specs.

Determines what slides should be created and in what order, before generation.
"""

import os
import re
from pptx import Presentation
from docx import Document

from bible_fetcher import fetch_verses, group_verses_for_slides
from slide_finder import find_slide, find_consecutive


def plan_slides(template, libraries, agenda, input_files, skip_intro=False, bible_page=None):
    """
    Build the ordered list of slide specs for the service.

    Args:
        template: Template Presentation (for fixed content lookup)
        libraries: List of library Presentations
        agenda: Parsed agenda dict with worship_order, sermon_outline, announcements
        input_files: Dict of filename → filepath
        skip_intro: If True, don't add intro slides (they come from intro PPTX)
        bible_page: Page number for scripture title (user input)

    Returns:
        List of slide spec dicts describing what to create
    """
    order = agenda.get("worship_order", [])
    slides = []

    # === Fixed pre-worship slides ===
    if not skip_intro:
        # Intro: copy slides up to (but not including) call_to_worship
        call_to_worship_idx = find_slide(template, "call_to_worship")
        if call_to_worship_idx is None:
            call_to_worship_idx = 6  # fallback

        for i in range(call_to_worship_idx):
            slides.append({"type": "copy_template", "prs": template, "index": i})

    # === Worship order items ===
    for item in order:
        itype = item.get("type")

        if itype == "call_to_worship":
            idx = find_slide(template, "call_to_worship")
            if idx is not None:
                slides.append({"type": "copy_template", "prs": template, "index": idx})
            slides.append({"type": "blank"})

        elif itype == "hymn":
            hymn_slides = _get_hymn_slides(libraries, item, input_files)
            slides.extend(hymn_slides)
            slides.append({"type": "blank"})

        elif itype == "prayer":
            for key in ["prayer", "lords_prayer_1", "lords_prayer_2"]:
                idx = find_slide(template, key)
                if idx is not None:
                    slides.append({"type": "copy_template", "prs": template, "index": idx})

        elif itype == "creed":
            for key in ["creed_1", "creed_2"]:
                idx = find_slide(template, key)
                if idx is not None:
                    slides.append({"type": "copy_template", "prs": template, "index": idx})
            slides.append({"type": "blank"})

        elif itype == "responsive":
            reading_slides = _get_reading_slides(libraries, item, input_files)
            slides.extend(reading_slides)
            slides.append({"type": "blank"})

        elif itype == "anthem":
            anthem_slides = _get_anthem_slides(libraries, item, input_files)
            slides.extend(anthem_slides)
            slides.append({"type": "blank"})

        elif itype == "scripture":
            scripture_slides = _get_scripture_slides(libraries, item, bible_page)
            slides.extend(scripture_slides)
            # No blank between scripture and sermon

        elif itype == "sermon":
            sermon_slides = _get_sermon_slides(agenda)
            slides.extend(sermon_slides)
            slides.append({"type": "blank"})

        elif itype == "offering":
            for key in ["offering_1", "offering_2"]:
                idx = find_slide(template, key)
                if idx is not None:
                    slides.append({"type": "copy_template", "prs": template, "index": idx})
            slides.append({"type": "blank"})

        elif itype == "communion":
            # Only include if 聖餐 is in the agenda
            for key in ["communion_1", "communion_2", "communion_3"]:
                idx = find_slide(template, key)
                if idx is not None:
                    slides.append({"type": "copy_template", "prs": template, "index": idx})
            slides.append({"type": "blank"})

        elif itype == "announcements":
            ann_slides = _get_announcement_slides(agenda)
            if ann_slides:
                idx = find_slide(template, "announce_title")
                if idx is not None:
                    slides.append({"type": "copy_template", "prs": template, "index": idx})
                slides.extend(ann_slides)
                slides.append({"type": "blank"})

        elif itype == "doxology":
            # Skip—handled in closing section below
            pass

        elif itype == "benediction":
            # Skip—handled in closing section below
            pass

    # === Fixed closing slides ===
    for key in ["doxology", "benediction", "quiet", "website"]:
        idx = find_slide(template, key)
        if idx is not None:
            slides.append({"type": "copy_template", "prs": template, "index": idx})

    return slides


def _get_hymn_slides(libraries, item, input_files, is_doxology=False):
    """Get slides for a hymn item, from input file or slide library."""
    num = item.get("number")
    title = item.get("title", "")

    # Try readable input file (PPTX only)
    matched = _match_file(num, title, input_files, [".pptx"])
    if matched:
        try:
            src = Presentation(matched)
            if src.slides:
                return [{"type": "copy_external", "prs": src, "index": i}
                        for i in range(len(src.slides))]
        except Exception:
            pass

    # Search all libraries
    for lib in libraries:
        indices = _find_hymn_slides_in_library(lib, num, title)
        if indices:
            return [{"type": "copy_external", "prs": lib, "index": i} for i in indices]

    label = f"{'頌榮' if is_doxology else '聖詩'} {num}: {title}" if num else title
    return [{"type": "hymn_placeholder", "label": label}]


def _get_reading_slides(libraries, item, input_files):
    """Get slides for the responsive reading."""
    num = item.get("number")
    title = item.get("title", "")

    matched = _match_file(num, title, input_files, [".pptx"])
    if matched:
        try:
            src = Presentation(matched)
            if src.slides:
                return [{"type": "copy_external", "prs": src, "index": i}
                        for i in range(len(src.slides))]
        except Exception:
            pass

    for lib in libraries:
        indices = _find_reading_slides_in_library(lib, num, title)
        if indices:
            return [{"type": "copy_external", "prs": lib, "index": i} for i in indices]

    return [{"type": "hymn_placeholder", "label": f"啟應文 {num}: {title}"}]


def _get_anthem_slides(libraries, item, input_files):
    """Get anthem title + lyrics slides from DOCX or library."""
    title = item.get("title", "")
    slides = [{"type": "anthem_title", "title": title}]

    # Try DOCX for lyrics
    matched = _match_file(None, title, input_files, [".docx", ".doc"])
    if matched:
        try:
            lyrics = _parse_docx_lyrics(matched)
            # Skip only the first title paragraph; keep duplicates as lyrics
            clean_title = re.sub(r"\s+", "", title)
            verses = []
            title_seen = False
            for v in lyrics:
                if not title_seen and re.sub(r"\s+", "", v) == clean_title:
                    title_seen = True
                    continue
                verses.append(v)
            if verses:
                for chunk in _group_anthem_verses(verses, max_lines=6):
                    slides.append({"type": "anthem_lyrics", "title": title, "lyrics": chunk})
                return slides
        except Exception:
            pass

    # Search libraries for anthem slides
    for lib in libraries:
        indices = _find_anthem_slides_in_library(lib, title)
        if indices:
            content_indices = [i for i in indices if i != indices[0]]
            if not content_indices:
                content_indices = indices
            return [{"type": "anthem_title", "title": title}] + \
                   [{"type": "copy_external", "prs": lib, "index": i} for i in content_indices]

    return slides


def _group_anthem_verses(verses, max_lines=6):
    """
    Pack verses into slide-sized chunks of exactly max_lines lines.

    Flattens paragraph boundaries — splits across paragraphs when needed
    so each slide gets a full max_lines count (last slide may have fewer).
    """
    all_lines = []
    for v in verses:
        all_lines.extend(v.split("\n"))

    groups = []
    for i in range(0, len(all_lines), max_lines):
        groups.append("\n".join(all_lines[i:i + max_lines]))
    return groups


def _get_scripture_slides(libraries, item, bible_page=None):
    """Get scripture slides from bible-api, libraries, or as title-only fallback."""
    ref = item.get("title", "")

    # Try fetching from 和合本 online
    verses = fetch_verses(ref)
    if verses:
        slides = [{"type": "scripture_title", "item": item, "bible_page": bible_page}]
        for group in group_verses_for_slides(verses):
            slides.append({"type": "scripture_verses", "ref": ref, "verses": group})
        return slides

    # Search libraries
    for lib in libraries:
        indices = _find_scripture_slides_in_library(lib, ref)
        if indices and len(indices) > 1:
            return [{"type": "copy_external", "prs": lib, "index": i} for i in indices]

    return [{"type": "scripture_title", "item": item, "bible_page": bible_page}]


def _get_sermon_slides(agenda):
    """Build sermon outline slides from the parsed agenda."""
    outline = agenda.get("sermon_outline", {})
    title = outline.get("title", "")
    scripture = outline.get("scripture", "")
    main_points = outline.get("main_points", [])
    preacher = ""

    # Get preacher from worship order
    for item in agenda.get("worship_order", []):
        if item.get("type") == "sermon":
            preacher = item.get("presenter", "")
            break

    slides = []

    # Title slide with sermon title + preacher
    slides.append({
        "type": "sermon_title",
        "title": title,
        "preacher": preacher,
        "scripture": f"《{scripture}》" if scripture else "",
    })

    # One slide per main point
    for mp in main_points:
        heading = mp.get("heading", "")
        points = mp.get("points", [])
        groups = _group_sermon_points(heading, points)
        for i, group in enumerate(groups):
            slides.append({
                "type": "sermon_point",
                "heading": heading,
                "points": group,
                "continuation": i > 0,
            })

    return slides


def _group_sermon_points(heading, points, chars_per_line=22, max_lines=10):
    """Split points into groups that fit within a single slide."""
    heading_lines = max(1, -(-len(heading) // chars_per_line)) if heading else 0
    reserved = 1 + heading_lines
    budget = max(1, max_lines - reserved)

    groups = []
    current = []
    current_lines = 0
    for pt in points:
        pt_lines = max(1, -(-len(pt) // chars_per_line))
        if current and current_lines + pt_lines > budget:
            groups.append(current)
            current = [pt]
            current_lines = pt_lines
        else:
            current.append(pt)
            current_lines += pt_lines
    if current:
        groups.append(current)
    return groups or [[]]


def _get_announcement_slides(agenda):
    """Build announcement slides — 1 item per slide."""
    announcements = agenda.get("announcements", {})
    slides = []

    for section, items in announcements.items():
        section_label = f"報告： {section}"
        for item in items:
            slides.append({
                "type": "announcement",
                "section": section_label,
                "text": item,
            })

    return slides


def _match_file(number, title, input_files, extensions):
    """Find an input file matching a hymn number or title."""
    for fname, fpath in input_files.items():
        ext = os.path.splitext(fname)[1].lower()
        if ext not in extensions:
            continue
        base = os.path.splitext(fname)[0]

        # Match by number prefix
        if number is not None:
            patterns = [
                rf"^0*{number}[-_\s]",
                rf"^0*{number}$",
            ]
            for pat in patterns:
                if re.match(pat, base, re.IGNORECASE):
                    return fpath

        # Match by title substring
        if title:
            clean_title = re.sub(r"\s+", "", title)
            clean_base = re.sub(r"\s+", "", base)
            if clean_title in clean_base or clean_base in clean_title:
                return fpath

    return None


def _find_hymn_slides_in_library(prs, number, title):
    """Find hymn slides in a presentation by number prefix or title."""
    seen = set()
    for i, slide in enumerate(prs.slides):
        text = _slide_text(slide)
        if number and re.search(rf"{number}\s*[：:]", text):
            seen.add(i)
            continue
        if title and number is None:
            clean_title = re.sub(r"\s+", "", title)
            if clean_title in re.sub(r"\s+", "", text):
                seen.add(i)

    # Fall back to title search if number search found nothing
    if not seen and title:
        clean_title = re.sub(r"\s+", "", title)
        for i, slide in enumerate(prs.slides):
            if clean_title in re.sub(r"\s+", "", _slide_text(slide)):
                seen.add(i)

    return sorted(seen) or None


def _find_reading_slides_in_library(prs, number, title):
    """Find responsive reading slides in a presentation."""
    seen = set()
    for i, slide in enumerate(prs.slides):
        text = _slide_text(slide)
        if number and re.search(rf"啟應文\s*{number}", text):
            seen.add(i)
            continue
        if title and not number:
            if re.sub(r"\s+", "", title) in re.sub(r"\s+", "", text):
                seen.add(i)
    return sorted(seen) or None


def _find_anthem_slides_in_library(prs, title):
    """Find anthem slides in a presentation by title text."""
    seen = set()
    clean = re.sub(r"\s+", "", title)
    for i, slide in enumerate(prs.slides):
        if clean in re.sub(r"\s+", "", _slide_text(slide)):
            seen.add(i)
    return sorted(seen) or None


def _find_scripture_slides_in_library(prs, ref):
    """Find scripture slides by book name / reference."""
    seen = set()
    book = re.split(r"\s*\d", ref)[0].strip()
    clean_book = re.sub(r"\s+", "", book)
    if len(clean_book) < 2:
        return None
    for i, slide in enumerate(prs.slides):
        clean_text = re.sub(r"\s+", "", _slide_text(slide))
        if clean_book in clean_text:
            seen.add(i)
    return sorted(seen) or None


def _slide_text(slide):
    """Extract all text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return " ".join(texts)


def _parse_docx_lyrics(docx_path):
    """Parse a DOCX lyrics file and return list of verse strings."""
    doc = Document(docx_path)
    verses = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            if lines:
                verses.append("\n".join(lines))
    return verses
