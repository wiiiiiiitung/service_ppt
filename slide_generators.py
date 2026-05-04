"""
Dynamic slide generation: create slides with content.

All styles/positions/fonts are read from service_config.SLIDE_STYLES.
No hardcoded formatting — everything drives from the config.
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

from service_config import SLIDE_STYLES
from styles import (
    get_style, apply_run_style, apply_paragraph_style, wrap_chinese_text,
    CYAN_BRIGHT, CYAN_SOFT, WHITE,
)
from slide_copier import get_layout
from bible_fetcher import get_testament


def add_placeholder_slide(out_prs, label):
    """Add a simple placeholder slide for missing content."""
    layout = get_layout(out_prs, "Blank")
    slide = out_prs.slides.add_slide(layout)
    tf = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.33), Inches(2)
    ).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2
    run = p.add_run()
    run.text = f"[{label}]"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = CYAN_BRIGHT


def add_anthem_title_slide(out_prs, title):
    """Anthem title slide: '獻詩: <title>' centered."""
    style = SLIDE_STYLES.get("anthem_title", {})
    layout = get_layout(out_prs, style.get("layout", "Blank"))
    slide = out_prs.slides.add_slide(layout)

    pos = style.get("pos", (Emu(981777), Emu(2786743)))
    size = style.get("size", (Emu(10501161), Emu(1015663)))

    tb = slide.shapes.add_textbox(pos[0], pos[1], size[0], size[1])
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = style.get("text", "獻詩: {title}").format(title=title)
    run.font.name = style.get("font", "標楷體")
    run.font.size = style.get("size_pt", Pt(60))
    run.font.bold = style.get("bold", True)
    if style.get("color"):
        run.font.color.rgb = style["color"]


def add_scripture_title_slide(out_prs, item, bible_page=None):
    """Scripture title slide with multi-run formatting."""
    style = SLIDE_STYLES.get("scripture_title", {})
    layout = get_layout(out_prs, style.get("layout", "Blank"))
    slide = out_prs.slides.add_slide(layout)

    pos = style.get("pos", (Emu(154746), Emu(2086708)))
    size = style.get("size", (Emu(12037254), Emu(2000250)))

    tb = slide.shapes.add_textbox(pos[0], pos[1], size[0], size[1])
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    ref = item.get("title", "")
    testament = get_testament(ref.split()[0]) if ref else "新約"

    # Build multi-run paragraph from config
    runs_config = style.get("runs", [])
    for i, run_cfg in enumerate(runs_config):
        if i > 0:
            run = p.add_run()
        else:
            run = p.add_run()

        text = run_cfg.get("text", "")
        if "{ref}" in text:
            text = text.format(ref=ref)
        if "{testament}" in text:
            text = text.format(testament=testament)
        if "{page}" in text:
            text = text.format(page=bible_page or "")

        run.text = text
        if run_cfg.get("font"):
            run.font.name = run_cfg["font"]
        if run_cfg.get("size_pt"):
            run.font.size = run_cfg["size_pt"]
        if run_cfg.get("bold"):
            run.font.bold = run_cfg["bold"]


def add_scripture_verse_slide(out_prs, ref, verses):
    """Scripture verse slide with ref bar + verse body."""
    layout = get_layout(out_prs, "Blank")
    slide = out_prs.slides.add_slide(layout)

    # Reference bar
    bar_style = SLIDE_STYLES.get("scripture_verse_bar", {})
    bpos = bar_style.get("pos", (Emu(0), Emu(0)))
    bsize = bar_style.get("size", (Emu(12192000), Emu(1561514)))

    tb_ref = slide.shapes.add_textbox(bpos[0], bpos[1], bsize[0], bsize[1])
    tf_ref = tb_ref.text_frame
    p_ref = tf_ref.paragraphs[0]
    p_ref.alignment = 1  # left
    run_ref = p_ref.add_run()
    run_ref.text = ref
    run_ref.font.name = bar_style.get("font", "標楷體")
    run_ref.font.size = bar_style.get("size_pt", Pt(40))
    run_ref.font.bold = bar_style.get("bold", True)

    # Verse body
    body_style = SLIDE_STYLES.get("scripture_verse_body", {})
    dpos = body_style.get("pos", (Emu(-22035), Emu(815926)))
    dsize = body_style.get("size", (Emu(12537195), Emu(6084277)))
    chars_per_line = body_style.get("chars_per_line", 18)

    tb_body = slide.shapes.add_textbox(dpos[0], dpos[1], dsize[0], dsize[1])
    tf_body = tb_body.text_frame
    tf_body.word_wrap = True

    first = True
    for v in verses:
        prefix = f"{v['verse']}."
        indent = " " * len(prefix)
        line = prefix + v['text']
        if chars_per_line:
            line = wrap_chinese_text(line, chars_per_line, indent=indent)

        p = tf_body.paragraphs[0] if first else tf_body.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.name = body_style.get("font", "DFKai-SB")
        run.font.size = body_style.get("size_pt", Pt(54))
        run.font.bold = body_style.get("bold", True)


def add_lyrics_slide(out_prs, title, lyrics):
    """Hymn/anthem lyrics slide using 詩歌 layout placeholders."""
    layout = get_layout(out_prs, "詩歌")
    slide = out_prs.slides.add_slide(layout)

    ph_list = list(slide.placeholders)
    lines = lyrics.split("\n")

    if len(ph_list) >= 2:
        ph_list[0].text = title
        tf = ph_list[1].text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
    else:
        tb = slide.shapes.add_textbox(Emu(0), Emu(126609), Emu(12192000), Emu(436100))
        tb.text_frame.paragraphs[0].text = title

        tb2 = slide.shapes.add_textbox(Emu(0), Emu(745435), Emu(12192000), Emu(6112565))
        tf = tb2.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line


def add_sermon_title_slide(out_prs, spec):
    """Sermon title slide: '今日信息' + title + preacher."""
    layout = get_layout(out_prs, "詩歌")
    slide = out_prs.slides.add_slide(layout)
    ph_list = list(slide.placeholders)

    sermon_title = spec.get("title", "")
    preacher = spec.get("preacher", "")
    style = SLIDE_STYLES.get("sermon_title", {})

    if len(ph_list) >= 2:
        # Title placeholder
        tf0 = ph_list[0].text_frame
        tf0.clear()
        p0 = tf0.paragraphs[0]
        run0 = p0.add_run()
        run0.text = style.get("header_text", "今日信息")
        run0.font.name = "標楷體"
        run0.font.bold = True

        # Body placeholder
        tf = ph_list[1].text_frame
        tf.clear()
        p_title = tf.paragraphs[0]

        # Sermon title
        run_t = p_title.add_run()
        run_t.text = sermon_title
        run_t.font.size = style.get("title_size_pt", Pt(80))

        # Spacer
        if preacher:
            run_sp = p_title.add_run()
            run_sp.text = "             "
            run_sp.font.name = style.get("spacer_font", "DFKai-SB")
            run_sp.font.size = style.get("spacer_size_pt", Pt(50))
            run_sp.font.bold = style.get("spacer_bold", True)

            # Preacher
            run_p = p_title.add_run()
            run_p.text = preacher
            run_p.font.name = style.get("preacher_font", "標楷體")
            run_p.font.size = style.get("preacher_size_pt", Pt(44))
            run_p.font.bold = style.get("preacher_bold", True)


def add_sermon_point_slide(out_prs, spec):
    """Sermon point slide: heading + bullet points."""
    layout = get_layout(out_prs, "詩歌")
    slide = out_prs.slides.add_slide(layout)

    heading = spec.get("heading", "")
    points = spec.get("points", [])
    if spec.get("continuation") and heading and "(續)" not in heading:
        heading = f"{heading} (續)"

    # Remove title placeholder
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 0:
            sp = ph._element
            sp.getparent().remove(sp)

    # Use body placeholder
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break

    style = SLIDE_STYLES.get("sermon_point", {})
    hdr_font = style.get("header_font", "標楷體")
    hdr_size = style.get("header_size_pt", Pt(44))

    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()

        p0 = tf.paragraphs[0]
        run0 = p0.add_run()
        run0.text = style.get("header_text", "今日信息")
        run0.font.name = hdr_font
        run0.font.size = hdr_size
        run0.font.bold = True

        p_heading = tf.add_paragraph()
        p_heading.alignment = 0  # left
        run_heading = p_heading.add_run()
        run_heading.text = heading

        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.alignment = 0
            run_pt = p_pt.add_run()
            run_pt.text = pt
    else:
        tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(12192000), Emu(6858000))
        tf = tb.text_frame
        p0 = tf.paragraphs[0]
        run0 = p0.add_run()
        run0.text = "今日信息"
        run0.font.name = hdr_font
        run0.font.size = hdr_size
        run0.font.bold = True
        for text in [heading] + points:
            p = tf.add_paragraph()
            p.alignment = 0
            run = p.add_run()
            run.text = text


def add_announcement_slide(out_prs, spec):
    """Announcement slide: one item per slide."""
    layout = get_layout(out_prs, "詩歌")
    slide = out_prs.slides.add_slide(layout)
    ph_list = list(slide.placeholders)

    section = spec.get("section", "報告")
    text = spec.get("text", "")

    if len(ph_list) >= 2:
        ph_list[0].text = section
        tf = ph_list[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = 0  # left
    else:
        tb = slide.shapes.add_textbox(
            Emu(998220), Emu(-152400), Emu(9403080), Emu(1143000)
        )
        tb.text_frame.paragraphs[0].text = section

        tb2 = slide.shapes.add_textbox(
            Emu(0), Emu(990600), Emu(12355033), Emu(5360964)
        )
        tf = tb2.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = 0  # left
        run = p.add_run()
        run.text = text
        run.font.size = Pt(18)
